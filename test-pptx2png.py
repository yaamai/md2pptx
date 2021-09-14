from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE


pres = Presentation("base.pptx")
print(pres.core_properties.__dict__)
print(pres.slide_master)
dir(pres.slide_master)
print(pres.slide_master.slide_layouts)
print(len(pres.slide_master.slide_layouts))
print(pres.slide_master.slide_layouts[0].shapes)
print(len(pres.slide_master.slide_layouts[0].shapes))

for layout in pres.slide_master.slide_layouts:
    print(len(layout.shapes))
    for shape in layout.shapes:
        print(shape)
        if hasattr(shape, "type"):
            print(shape.type)
        print(shape.left, shape.top, shape.width, shape.height)


print("---")
print(pres.slide_width, pres.slide_height)

print("---")
for shape in pres.slide_master.shapes:
    print(shape)
    if hasattr(shape, "type"):
        print(shape.type)
    print(shape.left, shape.top, shape.width, shape.height)

print("---")
for placeholder in pres.slide_master.placeholders:
    print(placeholder)
    print(placeholder.left, placeholder.top, placeholder.width, placeholder.height)


import drawSvg as draw

print("---")
print(pres.slide_width, pres.slide_height)

# drawSvg was designed for the coordinates x and y to increase rightward and upward. (#11)
FACTOR = 720
d = draw.Drawing(pres.slide_width/FACTOR, pres.slide_height/FACTOR)
g = draw.Group(transform="scale(1,-1) translate(0,{})".format(pres.slide_height/FACTOR))

for shape in pres.slide_master.shapes:
    print(shape.shape_type, shape.left, shape.top, shape.width, shape.height)
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        if shape.auto_shape_type == MSO_SHAPE.OVAL:
            r = shape.width/FACTOR
            c = shape.fill.fore_color.rgb
            line_color = shape.line.fill.fore_color.rgb
            line_width = shape.line.width/FACTOR
            g.append(draw.Circle(shape.left/FACTOR+r/2, shape.top/FACTOR+r/2, r/2, fill='#{}'.format(c), stroke='#{}'.format(line_color), stroke_width=line_width))
    elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
        line_color = shape.line.fill.fore_color.rgb
        line_width = shape.line.width/FACTOR
        g.append(draw.Line(shape.begin_x/FACTOR, shape.begin_y/FACTOR, shape.end_x/FACTOR, shape.end_y/FACTOR, stroke='#{}'.format(line_color), stroke_width=line_width))
    elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        print(shape.text)
        g.append(draw.Rectangle(shape.left/FACTOR, shape.top/FACTOR, shape.width/FACTOR, shape.height/FACTOR, fill='none', stroke='#1248ff', stroke_width=1000/FACTOR))
        g.append(draw.Text(shape.text, 200, 0, 0, text_anchor='middle', transform='scale(1,-1) translate({},{})'.format(shape.left/FACTOR, shape.top/FACTOR)))
    else:
        g.append(draw.Rectangle(shape.left/FACTOR, shape.top/FACTOR, shape.width/FACTOR, shape.height/FACTOR, fill='#1248ff'))
    # d.append(draw.Rectangle(shape.left/FACTOR, pres.slide_height/FACTOR-shape.top/FACTOR, shape.width/FACTOR, shape.height/FACTOR, fill='#1248ff'))
    # g.append(draw.Text('Text at (2,1)',0.5, 0,0, text_anchor='middle', transform='scale(1,-1) translate(2,1)'))

d.append(g)
d.setRenderSize(h=600)
d.saveSvg('out.svg')

