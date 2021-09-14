import io
from marko import Markdown
from marko import Parser
from marko.html_renderer import HTMLRenderer
from marko.renderer import Renderer
from pptx import Presentation

"""
- heading(level=1) starts new slide
- horizonal line starts new slide
- heading(level>1) as
"""

class PptxRenderer(Renderer):

    def render_document(self, element):
        self.pres = Presentation("base.pptx")
        self.render_children(element)

        return self.pres

    def render_raw_text(self, element):
        print(element)
        print(self.slide.shapes[1].text_frame.paragraphs)

        # text frame has one paragraph initially (with or w/o clear())
        if len(self.slide.shapes[1].text_frame.paragraphs) == 1:
            p = self.slide.shapes[1].text_frame.paragraphs[0]
        else:
            p = self.slide.shapes[1].text_frame.add_paragraph()

        p._pPr._add_buNone()
        p.text = element.children
        print(self.slide.shapes[1].text_frame.paragraphs)
        # import pdb; pdb.set_trace()
        # FIXME: rawtext in slide content placeholder appears in paragraph


    def render_heading(self, element):
        print(element.level)
        if element.level == 1:
            SLD_LAYOUT_TITLE_AND_CONTENT = 2
            slide_layout = self.pres.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
            self.slide = self.pres.slides.add_slide(slide_layout)
            self.slide.shapes.title.text = element.children[0].children
            for shape in self.slide.placeholders:
                print('%d %s %d' % (shape.placeholder_format.idx, shape.name, shape.placeholder_format.type))
            self.slide.shapes[1].text_frame.clear()
            print(self.slide.shapes[1].text_frame.paragraphs)
            return

        if element.level == 2:
            shape = self.slide.shapes[1]
            shape.text = element.children[0].children

    def render_children(self, element):
        print(element)
        if hasattr(element, "children") and isinstance(element.children, list):
            [self.render(child) for child in element.children]

def main():
    md = Markdown(parser=Parser, renderer=PptxRenderer)
    md_str = """
# slide title #1
slide body text
slide body text
slide body text
slide body text

# slide title #2
## slide body paragraph #1
## slide body paragraph #2
"""
    pres = md(md_str)

    with open("out.pptx", "wb") as f:
        pres.save(f)

if __name__ == '__main__':
    main()

