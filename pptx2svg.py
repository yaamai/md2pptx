from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE
import base64
import svgwrite

FACTOR = 360

"""
<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvGrpSpPr>
    <p:cNvPr id="361484" name="Group 12"/>
    <p:cNvGrpSpPr>
      <a:grpSpLocks noChangeAspect="1"/>
    </p:cNvGrpSpPr>
    <p:nvPr userDrawn="1"/>
  </p:nvGrpSpPr>
  <p:grpSpPr bwMode="auto">
    <a:xfrm>
      <a:off x="7607300" y="114300"/>
      <a:ext cx="1417638" cy="792163"/>
      <a:chOff x="4792" y="72"/>
      <a:chExt cx="893" cy="499"/>
    </a:xfrm>
  </p:grpSpPr>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr id="361483" name="AutoShape 11"/>
      <p:cNvSpPr>
        <a:spLocks noChangeAspect="1" noChangeArrowheads="1" noTextEdit="1"/>
      </p:cNvSpPr>
      <p:nvPr userDrawn="1"/>
    </p:nvSpPr>
    <p:spPr bwMode="gray">
      <a:xfrm>
        <a:off x="4792" y="72"/>
        <a:ext cx="893" cy="499"/>
      </a:xfrm>
      <a:prstGeom prst="rect">
        <a:avLst/>
      </a:prstGeom>
      <a:noFill/>
      <a:ln>
        <a:noFill/>
      </a:ln>
      <a:extLst>
        <a:ext uri="{909E8E84-426E-40DD-AFC4-6F175D3DCCD1}">
          <a14:hiddenFill xmlns:a14="http://schemas.microsoft.com/office/drawing/2010/main">
            <a:solidFill>
              <a:srgbClr val="FFFFFF"/>
            </a:solidFill>
          </a14:hiddenFill>
        </a:ext>
        <a:ext uri="{91240B29-F687-4F45-9708-019B960494DF}">
          <a14:hiddenLine xmlns:a14="http://schemas.microsoft.com/office/drawing/2010/main" w="9525">
            <a:solidFill>
              <a:srgbClr val="000000"/>
            </a:solidFill>
            <a:miter lim="800000"/>
            <a:headEnd/>
            <a:tailEnd/>
          </a14:hiddenLine>
        </a:ext>
      </a:extLst>
    </p:spPr>
    <p:txBody>
      <a:bodyPr/>
      <a:lstStyle/>
      <a:p>
        <a:endParaRPr lang="ja-JP" altLang="en-US"/>
      </a:p>
    </p:txBody>
  </p:sp>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr id="361485" name="Freeform 13"/>
      <p:cNvSpPr>
        <a:spLocks/>
      </p:cNvSpPr>
      <p:nvPr userDrawn="1"/>
    </p:nvSpPr>
    <p:spPr bwMode="gray">
      <a:xfrm>
        <a:off x="5232" y="157"/>
        <a:ext cx="113" cy="88"/>
      </a:xfrm>
      <a:custGeom>
        <a:avLst/>
        <a:gdLst>
          <a:gd name="T0" fmla="*/ 961 w 2365"/>
          <a:gd name="T1" fmla="*/ 764 h 1825"/>
          <a:gd name="T2" fmla="*/ 599 w 2365"/>
          <a:gd name="T3" fmla="*/ 642 h 1825"/>
          <a:gd name="T4" fmla="*/ 1 w 2365"/>
          <a:gd name="T5" fmla="*/ 1235 h 1825"/>
          <a:gd name="T6" fmla="*/ 599 w 2365"/>
          <a:gd name="T7" fmla="*/ 1825 h 1825"/>
          <a:gd name="T8" fmla="*/ 1053 w 2365"/>
          <a:gd name="T9" fmla="*/ 1619 h 1825"/>
          <a:gd name="T10" fmla="*/ 1053 w 2365"/>
          <a:gd name="T11" fmla="*/ 1293 h 1825"/>
          <a:gd name="T12" fmla="*/ 771 w 2365"/>
          <a:gd name="T13" fmla="*/ 1568 h 1825"/>
          <a:gd name="T14" fmla="*/ 599 w 2365"/>
          <a:gd name="T15" fmla="*/ 1604 h 1825"/>
          <a:gd name="T16" fmla="*/ 225 w 2365"/>
          <a:gd name="T17" fmla="*/ 1235 h 1825"/>
          <a:gd name="T18" fmla="*/ 599 w 2365"/>
          <a:gd name="T19" fmla="*/ 863 h 1825"/>
          <a:gd name="T20" fmla="*/ 861 w 2365"/>
          <a:gd name="T21" fmla="*/ 972 h 1825"/>
          <a:gd name="T22" fmla="*/ 1093 w 2365"/>
          <a:gd name="T23" fmla="*/ 1239 h 1825"/>
          <a:gd name="T24" fmla="*/ 1630 w 2365"/>
          <a:gd name="T25" fmla="*/ 1473 h 1825"/>
          <a:gd name="T26" fmla="*/ 2365 w 2365"/>
          <a:gd name="T27" fmla="*/ 741 h 1825"/>
          <a:gd name="T28" fmla="*/ 1630 w 2365"/>
          <a:gd name="T29" fmla="*/ 0 h 1825"/>
          <a:gd name="T30" fmla="*/ 1053 w 2365"/>
          <a:gd name="T31" fmla="*/ 295 h 1825"/>
          <a:gd name="T32" fmla="*/ 1053 w 2365"/>
          <a:gd name="T33" fmla="*/ 735 h 1825"/>
          <a:gd name="T34" fmla="*/ 1630 w 2365"/>
          <a:gd name="T35" fmla="*/ 232 h 1825"/>
          <a:gd name="T36" fmla="*/ 2134 w 2365"/>
          <a:gd name="T37" fmla="*/ 741 h 1825"/>
          <a:gd name="T38" fmla="*/ 1630 w 2365"/>
          <a:gd name="T39" fmla="*/ 1245 h 1825"/>
          <a:gd name="T40" fmla="*/ 1303 w 2365"/>
          <a:gd name="T41" fmla="*/ 1125 h 1825"/>
          <a:gd name="T42" fmla="*/ 961 w 2365"/>
          <a:gd name="T43" fmla="*/ 764 h 1825"/>
        </a:gdLst>
        <a:ahLst/>
        <a:cxnLst>
          <a:cxn ang="0">
            <a:pos x="T0" y="T1"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T2" y="T3"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T4" y="T5"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T6" y="T7"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T8" y="T9"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T10" y="T11"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T12" y="T13"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T14" y="T15"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T16" y="T17"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T18" y="T19"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T20" y="T21"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T22" y="T23"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T24" y="T25"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T26" y="T27"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T28" y="T29"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T30" y="T31"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T32" y="T33"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T34" y="T35"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T36" y="T37"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T38" y="T39"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T40" y="T41"/>
          </a:cxn>
          <a:cxn ang="0">
            <a:pos x="T42" y="T43"/>
          </a:cxn>
        </a:cxnLst>
        <a:rect l="0" t="0" r="r" b="b"/>
        <a:pathLst>
          <a:path w="2365" h="1825">
            <a:moveTo>
              <a:pt x="961" y="764"/>
            </a:moveTo>
            <a:cubicBezTo>
              <a:pt x="875" y="688"/>
              <a:pt x="736" y="643"/>
              <a:pt x="599" y="642"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="270" y="641"/>
              <a:pt x="2" y="900"/>
              <a:pt x="1" y="1235"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="0" y="1563"/>
              <a:pt x="270" y="1824"/>
              <a:pt x="599" y="1825"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="783" y="1825"/>
              <a:pt x="943" y="1751"/>
              <a:pt x="1053" y="1619"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="1053" y="1293"/>
              <a:pt x="1053" y="1293"/>
              <a:pt x="1053" y="1293"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="994" y="1394"/>
              <a:pt x="877" y="1524"/>
              <a:pt x="771" y="1568"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="717" y="1590"/>
              <a:pt x="662" y="1604"/>
              <a:pt x="599" y="1604"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="394" y="1604"/>
              <a:pt x="225" y="1445"/>
              <a:pt x="225" y="1235"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="225" y="1041"/>
              <a:pt x="380" y="862"/>
              <a:pt x="599" y="863"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="701" y="863"/>
              <a:pt x="793" y="905"/>
              <a:pt x="861" y="972"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="931" y="1040"/>
              <a:pt x="1041" y="1183"/>
              <a:pt x="1093" y="1239"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="1227" y="1383"/>
              <a:pt x="1419" y="1473"/>
              <a:pt x="1630" y="1473"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="2036" y="1474"/>
              <a:pt x="2365" y="1146"/>
              <a:pt x="2365" y="741"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="2365" y="336"/>
              <a:pt x="2035" y="0"/>
              <a:pt x="1630" y="0"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="1395" y="0"/>
              <a:pt x="1187" y="122"/>
              <a:pt x="1053" y="295"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="1053" y="735"/>
              <a:pt x="1053" y="735"/>
              <a:pt x="1053" y="735"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="1155" y="455"/>
              <a:pt x="1340" y="232"/>
              <a:pt x="1630" y="232"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="1909" y="232"/>
              <a:pt x="2135" y="463"/>
              <a:pt x="2134" y="741"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="2134" y="1020"/>
              <a:pt x="1909" y="1246"/>
              <a:pt x="1630" y="1245"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="1506" y="1244"/>
              <a:pt x="1391" y="1199"/>
              <a:pt x="1303" y="1125"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="1194" y="1040"/>
              <a:pt x="1074" y="860"/>
              <a:pt x="961" y="764"/>
            </a:cubicBezTo>
          </a:path>
        </a:pathLst>
      </a:custGeom>
      <a:solidFill>
        <a:srgbClr val="FF0000"/>
      </a:solidFill>
      <a:ln>
        <a:noFill/>
      </a:ln>
      <a:extLst>
        <a:ext uri="{91240B29-F687-4F45-9708-019B960494DF}">
          <a14:hiddenLine xmlns:a14="http://schemas.microsoft.com/office/drawing/2010/main" w="9525">
            <a:solidFill>
              <a:srgbClr val="000000"/>
            </a:solidFill>
            <a:round/>
            <a:headEnd/>
            <a:tailEnd/>
          </a14:hiddenLine>
        </a:ext>
      </a:extLst>
    </p:spPr>
    <p:txBody>
      <a:bodyPr/>
      <a:lstStyle/>
      <a:p>
        <a:endParaRPr lang="ja-JP" altLang="en-US"/>
      </a:p>
    </p:txBody>
  </p:sp>
"""
def _convert_shape(dwg, shape):
    print(shape.shape_type, shape, shape.left, shape.top, shape.width, shape.height)
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        # print(shape.element.xml)
        if shape.auto_shape_type == MSO_SHAPE.OVAL:
            pass
            # r = shape.width/FACTOR
            # # import pdb; pdb.set_trace()
            # # fill_color = 'rgb({}, {}, {})'.format(*shape.fill.fore_color.rgb)
            # fill_color = ''
            # line_color = 'rgb({}, {}, {})'.format(*shape.line.fill.fore_color.rgb)
            # line_width = shape.line.width/FACTOR
            # return dwg.circle((shape.left/FACTOR+r/2, shape.top/FACTOR+r/2), r/2, fill=fill_color, stroke=line_color, stroke_width=line_width)
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        # print(shape, shape.name, shape.element.xml, "SHAPE len = {}".format(len(shape.shapes)))
        r = dwg.rect(insert=(shape.left/FACTOR, shape.top/FACTOR),
                        size=(shape.width/FACTOR, shape.height/FACTOR), fill='blue', opacity=0.1, stroke='gray', stroke_width=10)
        dwg.add(r)
        group = dwg.g()
        group.translate(shape.left/FACTOR, shape.top/FACTOR)
        for s in shape.shapes:
            group.add(_convert_shape(dwg, s))
        return group
    elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
        pass
    elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
        # import pdb; pdb.set_trace()
        group = dwg.g()
        group.scale(shape.width/shape._parent.parent.element.xfrm.chExt.cx, shape.height/shape._parent.parent.element.xfrm.chExt.cy)
        for p in shape.element.spPr.custGeom.pathLst:
            # import pdb; pdb.set_trace()
            print(shape.line.width)
            line_width = shape.line.width/FACTOR
            if shape.line.width == 0:
                # SVG's default stroke width == 1
                line_width = 30
            color = 'red'
            # color = 'rgb({}, {}, {})'.format(*shape.line.color.rgb)
            path = dwg.path(fill="none", stroke=color, stroke_width=line_width)
            path.translate(shape.left/FACTOR, shape.top/FACTOR)
            path.scale((shape.width)/p.w, (shape.height)/p.h)
            # path.scale(shape._parent.parent.element.xfrm.chExt.cx/p.w, shape._parent.parent.element.xfrm.chExt.cy/p.h)
            # print("Scale: ", shape.left/FACTOR, shape.top/FACTOR, (shape.width/360)/p.w, (shape.height/360)/p.h, shape.width/FACTOR, shape.height/FACTOR, p.w, p.h)
            for command in p:
                if command.tag == '{http://schemas.openxmlformats.org/drawingml/2006/main}moveTo':
                    path.push('M', command.pt.x, command.pt.y)
                elif command.tag == '{http://schemas.openxmlformats.org/drawingml/2006/main}lnTo':
                    # import pdb; pdb.set_trace()
                    path.push('L', command.pt.x, command.pt.y)
                elif command.tag == '{http://schemas.openxmlformats.org/drawingml/2006/main}cubicBezTo':
                    path.push('C', command[0].x, command[0].y, command[1].x, command[1].y, command[2].x, command[2].y)
                elif command.tag == '{http://schemas.openxmlformats.org/drawingml/2006/main}close':
                    path.push('Z')
                else:
                    # pass
                    import pdb; pdb.set_trace()
            group.add(path)
        return group

    elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        return dwg.rect(insert=(shape.left/FACTOR, shape.top/FACTOR),
                        size=(shape.width/FACTOR, shape.height/FACTOR), fill='gray', opacity=0.1, stroke='gray', stroke_width=10)

    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        b64_data = base64.b64encode(shape.image.blob).decode()
        data_uri = 'data:{};base64,{}'.format('image/png', b64_data)
        return dwg.image(data_uri, insert=(shape.left/FACTOR, shape.top/FACTOR),
                         size=(shape.width/FACTOR, shape.height/FACTOR))

    return dwg.rect(insert=(shape.left/FACTOR, shape.top/FACTOR),
                    size=(shape.width/FACTOR, shape.height/FACTOR), fill='blue', stroke='red', stroke_width=10)


def main():
    pres = Presentation("fj2.pptx")
    print(pres.slide_width, pres.slide_height)
    dwg = svgwrite.Drawing(filename="out.svg", size=(1200, 1080), debug=True)
    dwg.viewbox(0, 0, pres.slide_width/FACTOR, pres.slide_height/FACTOR)

    for shape in pres.slide_master.shapes:
        dwg.add(_convert_shape(dwg, shape))

    dwg.save()

main()
